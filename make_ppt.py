from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ─────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan-blue
ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
GREEN     = RGBColor(0x06, 0xD6, 0xA0)
YELLOW    = RGBColor(0xFF, 0xD1, 0x66)
ORANGE    = RGBColor(0xF7, 0x7F, 0x00)
RED       = RGBColor(0xEF, 0x47, 0x6F)
BLOCK_BG  = RGBColor(0x1B, 0x2E, 0x44)   # card background

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════
def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background() if border is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def htxt(slide, text, l, t, w, h, size=Pt(20), bold=True, color=ACCENT, align=PP_ALIGN.LEFT):
    return txt(slide, text, l, t, w, h, size=size, bold=bold, color=color, align=align)

def add_image(slide, path, l, t, w, h):
    try:
        slide.shapes.add_picture(path, l, t, w, h)
    except Exception as e:
        print(f"  [image skipped] {path}: {e}")

def divider(slide, t, color=ACCENT):
    box(slide, Inches(0.5), t, Inches(12.33), Pt(2), fill=color)

def bullet_list(slide, items, l, t, w, h,
                size=Pt(16), color=WHITE, dot_color=ACCENT, indent=Inches(0.3)):
    line_h = size * 1.6
    y = t
    for item in items:
        # dot
        dot = slide.shapes.add_shape(9, l, y + size * 0.35, Pt(7), Pt(7))  # 9=oval
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        txt(slide, item, l + indent, y, w - indent, line_h, size=size, color=color)
        y += line_h

def card(slide, l, t, w, h, title, content_lines,
         title_color=ACCENT, content_size=Pt(14), title_size=Pt(16)):
    box(slide, l, t, w, h, fill=BLOCK_BG, border=ACCENT, border_w=Pt(1))
    txt(slide, title, l + Inches(0.15), t + Inches(0.1),
        w - Inches(0.3), Inches(0.4), size=title_size, bold=True, color=title_color)
    cy = t + Inches(0.55)
    lh = content_size * 1.5
    for line in content_lines:
        txt(slide, line, l + Inches(0.15), cy, w - Inches(0.3), lh,
            size=content_size, color=GRAY)
        cy += lh

def section_badge(slide, label, l=Inches(0.5), t=Inches(0.2)):
    box(slide, l, t, Inches(2.2), Inches(0.38), fill=ACCENT)
    txt(slide, label, l + Inches(0.12), t + Inches(0.05),
        Inches(2), Inches(0.3), size=Pt(14), bold=True, color=BG_DARK)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# top accent bar
box(s, 0, 0, prs.slide_width, Inches(0.08), fill=ACCENT)
box(s, 0, prs.slide_height - Inches(0.08), prs.slide_width, Inches(0.08), fill=ACCENT)

# decorative circle (large, faint)
c = s.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(6), Inches(6))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x00, 0x5F, 0x7A)
c.line.fill.background()

c2 = s.shapes.add_shape(9, Inches(10.5), Inches(-0.5), Inches(4), Inches(4))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x00, 0x3F, 0x5A)
c2.line.fill.background()

# FFA-Net label
txt(s, "FFA-Net", Inches(0.6), Inches(1.2), Inches(11), Inches(1.6),
    size=Pt(72), bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# subtitle
txt(s, "Feature Fusion Attention Network for Single Image Dehazing",
    Inches(0.6), Inches(2.8), Inches(10), Inches(1.0),
    size=Pt(26), bold=False, color=WHITE)

# conference badge
box(s, Inches(0.6), Inches(3.85), Inches(2.5), Inches(0.45), fill=ACCENT)
txt(s, "AAAI 2020", Inches(0.7), Inches(3.88), Inches(2.3), Inches(0.4),
    size=Pt(18), bold=True, color=BG_DARK)

# authors
txt(s, "Xu Qin, Zhilin Wang et al.",
    Inches(0.6), Inches(4.45), Inches(9), Inches(0.5),
    size=Pt(18), color=ACCENT2)
txt(s, "Peking University  ·  Beijing University of Aeronautics & Astronautics",
    Inches(0.6), Inches(4.95), Inches(9), Inches(0.5),
    size=Pt(15), color=GRAY, italic=True)

# bottom tag
txt(s, "arXiv: 1911.07559",
    Inches(0.6), Inches(6.5), Inches(4), Inches(0.4),
    size=Pt(13), color=GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — 目录
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "目录 / Contents")
txt(s, "报告大纲", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(1.4))

items = [
    ("01", "研究背景与问题定义"),
    ("02", "核心方法：FFA-Net 架构"),
    ("03", "注意力机制详解（CA / PA）"),
    ("04", "特征融合策略"),
    ("05", "训练流程与损失函数"),
    ("06", "数据集与实验设置"),
    ("07", "性能评估与对比结果"),
    ("08", "效果展示"),
    ("09", "代码结构"),
    ("10", "总结与展望"),
]
cols = 2
col_w = Inches(5.8)
lx = [Inches(0.6), Inches(7.0)]
ty = Inches(1.6)
row_h = Inches(0.52)
for i, (num, title) in enumerate(items):
    col = i % cols
    row = i // cols
    l = lx[col]
    t = ty + row * row_h
    box(s, l, t, col_w, row_h - Inches(0.06), fill=BLOCK_BG, border=ACCENT, border_w=Pt(1))
    txt(s, num, l + Inches(0.1), t + Inches(0.06), Inches(0.5), row_h - Inches(0.1),
        size=Pt(18), bold=True, color=ACCENT)
    txt(s, title, l + Inches(0.6), t + Inches(0.06), col_w - Inches(0.8), row_h - Inches(0.1),
        size=Pt(17), color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — 研究背景
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "01  研究背景")
txt(s, "单图像去雾：问题与挑战", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Haze model box
box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.3), fill=BLOCK_BG, border=YELLOW, border_w=Pt(1.5))
txt(s, "物理雾霾模型  I(x) = J(x)·t(x) + A·(1 − t(x))",
    Inches(0.7), Inches(1.7), Inches(11.8), Inches(0.6),
    size=Pt(21), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txt(s, "I: 有雾图像    J: 清晰场景    t(x): 透射图    A: 大气光照",
    Inches(0.7), Inches(2.2), Inches(11.8), Inches(0.5),
    size=Pt(15), color=GRAY, align=PP_ALIGN.CENTER)

cards_data = [
    ("雾霾的危害", ["降低图像能见度与对比度", "影响目标检测 / 分割 / 识别", "影响自动驾驶安全性"], ACCENT),
    ("现有方法局限", ["先验方法（DCP）假设强，易失效", "CNN 方法缺乏空间差异化处理", "不同位置雾浓度差异未被充分利用"], ORANGE),
    ("FFA-Net 的思路", ["引入像素级注意力（PALayer）", "引入通道级注意力（CALayer）", "多组特征跨组融合（Feature Fusion）"], GREEN),
]
cx = Inches(0.5)
cw = Inches(4.0)
ct = Inches(3.1)
ch = Inches(3.6)
for title, lines, tc in cards_data:
    card(s, cx, ct, cw, ch, title, lines, title_color=tc)
    cx += Inches(4.17)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — 整体架构
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "02  网络架构")
txt(s, "FFA-Net 整体结构", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Architecture diagram (text-art)
# Layers: Input -> PreConv -> G1 -> G2 -> G3 -> CrossGroupCA -> PALayer -> PostConv -> +Input -> Output
def arch_block(slide, label, sublabel, l, t, w=Inches(1.55), h=Inches(0.85), bc=BLOCK_BG, tc=ACCENT):
    box(slide, l, t, w, h, fill=bc, border=tc, border_w=Pt(1.5))
    txt(slide, label, l, t + Inches(0.06), w, Inches(0.44),
        size=Pt(14), bold=True, color=tc, align=PP_ALIGN.CENTER)
    if sublabel:
        txt(slide, sublabel, l, t + Inches(0.45), w, Inches(0.32),
            size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)

def arrow(slide, l, t, w=Inches(0.3)):
    box(slide, l, t + Inches(0.38), w, Pt(2), fill=ACCENT)
    # arrowhead
    tri = slide.shapes.add_shape(5, l + w - Pt(8), t + Inches(0.3),
                                  Pt(8), Inches(0.18))  # 5=right triangle
    tri.fill.solid(); tri.fill.fore_color.rgb = ACCENT
    tri.line.fill.background()

ty = Inches(1.65)
blocks = [
    ("Input\nHazy", "H×W×3", Inches(0.4)),
    ("PreConv", "3→64\n3×3", Inches(2.25)),
    ("Group 1\n(G1)", "N Blocks\nResidual", Inches(4.1)),
    ("Group 2\n(G2)", "N Blocks\nResidual", Inches(5.95)),
    ("Group 3\n(G3)", "N Blocks\nResidual", Inches(7.8)),
    ("Cross-Group\nCA Fusion", "Concat→\nWeighted Sum", Inches(9.65)),
    ("PALayer", "Pixel\nAttention", Inches(11.5)),
]
bw = Inches(1.55)
bh = Inches(0.95)
ar_w = Inches(0.2)
for i, (lbl, sub, lx) in enumerate(blocks):
    col = ACCENT if i not in (2,3,4) else (GREEN if i in (2,3,4) else ACCENT)
    if lbl.startswith("Cross"):
        col = YELLOW
    elif lbl == "PALayer":
        col = ORANGE
    elif "Input" in lbl:
        col = GRAY
    arch_block(s, lbl, sub, lx, ty, bw, bh, tc=col)
    if i < len(blocks)-1:
        arrow(s, lx + bw, ty, ar_w)

# PostConv + residual
box(s, Inches(0.4), ty + bh + Inches(0.4), Inches(13.0), Pt(2), fill=ACCENT)

# PostConv box
box(s, Inches(4.1), ty + bh + Inches(0.1), Inches(1.55), Inches(0.7),
    fill=BLOCK_BG, border=ACCENT2, border_w=Pt(1))
txt(s, "PostConv\n64→3", Inches(4.1), ty + bh + Inches(0.12), Inches(1.55), Inches(0.6),
    size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER)

# residual add
box(s, Inches(6.1), ty + bh + Inches(0.1), Inches(1.3), Inches(0.7),
    fill=BLOCK_BG, border=GREEN, border_w=Pt(1))
txt(s, "+ Input\n(Residual)", Inches(6.1), ty + bh + Inches(0.12), Inches(1.3), Inches(0.6),
    size=Pt(11), color=GREEN, align=PP_ALIGN.CENTER)

# Output label
box(s, Inches(7.8), ty + bh + Inches(0.1), Inches(1.3), Inches(0.7),
    fill=BLOCK_BG, border=WHITE, border_w=Pt(1))
txt(s, "Output\nClear Img", Inches(7.8), ty + bh + Inches(0.12), Inches(1.3), Inches(0.6),
    size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Key design notes
txt(s, "关键设计要点", Inches(0.5), Inches(4.2), Inches(5), Inches(0.4),
    size=Pt(16), bold=True, color=ACCENT)
notes = [
    "输入输出维度一致（H×W×3），通过残差学习去雾残差",
    "3 组 Group（G1/G2/G3）串行堆叠，每组含 N 个 Block（默认 N=19）",
    "三组输出跨组拼接后经 Channel Attention 加权融合",
    "最终通过 Pixel Attention 进行空间精细化",
    "整体参数量适中，适合 GPU 端到端训练",
]
bullet_list(s, notes, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.5),
            size=Pt(15), dot_color=ACCENT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Block 与 Group 结构
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "02  基本单元")
txt(s, "Block 与 Group 结构", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Block diagram (left half)
txt(s, "Block（基本块）", Inches(0.5), Inches(1.55), Inches(5.5), Inches(0.4),
    size=Pt(18), bold=True, color=GREEN)

bblocks = [
    ("Input x", ""),
    ("Conv 3×3 + ReLU", "dim→dim"),
    ("+ x  (skip)", ""),
    ("Conv 3×3", "dim→dim"),
    ("CALayer", "Channel Attn"),
    ("PALayer", "Pixel Attn"),
    ("+ x  (skip)", ""),
    ("Output", ""),
]
bl = Inches(0.7)
bt = Inches(2.0)
bw2 = Inches(2.4)
bh2 = Inches(0.48)
gap = Inches(0.06)
colors2 = [WHITE, ACCENT, GREEN, ACCENT, YELLOW, ORANGE, GREEN, WHITE]
for i, (lbl, sub) in enumerate(bblocks):
    col = colors2[i]
    fc = BLOCK_BG if col != WHITE else RGBColor(0x28,0x3F,0x58)
    box(s, bl, bt, bw2, bh2, fill=fc, border=col, border_w=Pt(1))
    txt(s, lbl, bl, bt + Inches(0.04), bw2, bh2 - Inches(0.04),
        size=Pt(13), bold=(col != GRAY), color=col, align=PP_ALIGN.CENTER)
    bt += bh2 + gap

# Group diagram (right half)
txt(s, "Group（组）", Inches(6.8), Inches(1.55), Inches(5.5), Inches(0.4),
    size=Pt(18), bold=True, color=ACCENT)

gblocks = [
    ("Input x", WHITE),
    ("Block 1", GREEN),
    ("Block 2", GREEN),
    ("  ···  (N Blocks)", GRAY),
    ("Block N", GREEN),
    ("Conv 3×3", ACCENT),
    ("+ x  (skip)", GREEN),
    ("Output", WHITE),
]
gl = Inches(7.0)
gt = Inches(2.0)
gw = Inches(2.4)
gh = Inches(0.48)
for lbl, col in gblocks:
    fc = BLOCK_BG if col != WHITE else RGBColor(0x28,0x3F,0x58)
    box(s, gl, gt, gw, gh, fill=fc, border=col, border_w=Pt(1))
    txt(s, lbl, gl, gt + Inches(0.04), gw, gh - Inches(0.04),
        size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)
    gt += gh + gap

# dividing line
box(s, Inches(6.3), Inches(1.5), Pt(2), Inches(5.8), fill=BLOCK_BG)

# Code snippet
txt(s, "关键代码（Block.forward）", Inches(3.5), Inches(1.55), Inches(3.5), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT2)
code = (
    "def forward(self, x):\n"
    "    res = self.act1(self.conv1(x))\n"
    "    res = res + x           # skip 1\n"
    "    res = self.conv2(res)\n"
    "    res = self.calayer(res) # channel attn\n"
    "    res = self.palayer(res) # pixel attn\n"
    "    res += x                # skip 2\n"
    "    return res"
)
box(s, Inches(3.5), Inches(2.0), Inches(3.2), Inches(2.5), fill=RGBColor(0x10,0x20,0x30))
txt(s, code, Inches(3.6), Inches(2.05), Inches(3.0), Inches(2.4),
    size=Pt(11), color=ACCENT2, italic=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Channel Attention
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "03  注意力机制")
txt(s, "通道注意力 — CALayer", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# formula box
box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9),
    fill=BLOCK_BG, border=YELLOW, border_w=Pt(1.5))
txt(s, "CALayer:  y = AvgPool(x)  →  Conv(C→C/8) + ReLU  →  Conv(C/8→C) + Sigmoid  →  x ⊗ y",
    Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.8),
    size=Pt(16), color=YELLOW, align=PP_ALIGN.CENTER)

# diagram boxes
ca_steps = [
    ("Feature Map\n[B,C,H,W]", GRAY, ""),
    ("Avg Pool\n→ [B,C,1,1]", ACCENT, "全局池化"),
    ("Conv C→C/8\n+ReLU", ACCENT, "压缩"),
    ("Conv C/8→C\n+Sigmoid", ACCENT, "激活"),
    ("× Input\n(逐通道相乘)", GREEN, "通道加权"),
    ("Output\n[B,C,H,W]", WHITE, ""),
]
cx0 = Inches(0.5)
ct0 = Inches(2.7)
cw0 = Inches(1.9)
ch0 = Inches(1.1)
for i, (lbl, col, sub) in enumerate(ca_steps):
    fc = BLOCK_BG if col not in (WHITE,) else RGBColor(0x28,0x3F,0x58)
    box(s, cx0, ct0, cw0, ch0, fill=fc, border=col, border_w=Pt(1.5))
    txt(s, lbl, cx0, ct0 + Inches(0.05), cw0, ch0 - Inches(0.05),
        size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)
    if sub:
        txt(s, sub, cx0, ct0 + ch0, cw0, Inches(0.3),
            size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)
    if i < len(ca_steps)-1:
        arrow(s, cx0 + cw0, ct0 + ch0 * 0.3, Inches(0.3))
    cx0 += cw0 + Inches(0.3)

txt(s, "作用：对不同特征通道赋予不同权重，强调信息量丰富的通道，抑制冗余通道",
    Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.5),
    size=Pt(16), color=ACCENT2, italic=True)

# code snippet
txt(s, "CALayer 核心代码", Inches(0.5), Inches(4.65), Inches(5), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT2)
ca_code = (
    "self.ca = nn.Sequential(\n"
    "    nn.AdaptiveAvgPool2d(1),\n"
    "    nn.Conv2d(C, C//8, 1),\n"
    "    nn.ReLU(inplace=True),\n"
    "    nn.Conv2d(C//8, C, 1),\n"
    "    nn.Sigmoid()\n"
    ")\n"
    "def forward(self, x):\n"
    "    return x * self.ca(x)"
)
box(s, Inches(0.5), Inches(5.1), Inches(5.5), Inches(2.1), fill=RGBColor(0x10,0x20,0x30))
txt(s, ca_code, Inches(0.65), Inches(5.15), Inches(5.3), Inches(2.0),
    size=Pt(12), color=ACCENT2, italic=True)

# why CA matters
txt(s, "为什么通道注意力重要？", Inches(6.5), Inches(4.65), Inches(6), Inches(0.4),
    size=Pt(15), bold=True, color=YELLOW)
ca_notes = [
    "雾霾对不同颜色通道影响程度不同（蓝通道受影响更大）",
    "压缩比 1/8 在感受野与参数量之间取得平衡",
    "Sigmoid 归一化保证权重在 [0,1] 区间",
    "全局平均池化可捕捉全图统计信息",
]
bullet_list(s, ca_notes, Inches(6.5), Inches(5.1), Inches(6.3), Inches(2.2),
            size=Pt(14), dot_color=YELLOW)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Pixel Attention
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "03  注意力机制")
txt(s, "像素注意力 — PALayer", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9),
    fill=BLOCK_BG, border=ORANGE, border_w=Pt(1.5))
txt(s, "PALayer:  y = Conv(C→C/8) + ReLU  →  Conv(C/8→1) + Sigmoid  →  x ⊗ y  (空间掩码)",
    Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.8),
    size=Pt(16), color=ORANGE, align=PP_ALIGN.CENTER)

pa_steps = [
    ("Feature Map\n[B,C,H,W]", GRAY, ""),
    ("Conv C→C/8\n+ReLU", ORANGE, "特征压缩"),
    ("Conv C/8→1\n+Sigmoid", ORANGE, "单通道掩码"),
    ("× Input\n(逐像素相乘)", GREEN, "空间加权"),
    ("Output\n[B,C,H,W]", WHITE, ""),
]
cx0 = Inches(0.9)
ct0 = Inches(2.7)
cw0 = Inches(2.0)
ch0 = Inches(1.1)
for i, (lbl, col, sub) in enumerate(pa_steps):
    fc = BLOCK_BG if col != WHITE else RGBColor(0x28,0x3F,0x58)
    box(s, cx0, ct0, cw0, ch0, fill=fc, border=col, border_w=Pt(1.5))
    txt(s, lbl, cx0, ct0 + Inches(0.05), cw0, ch0 - Inches(0.05),
        size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER)
    if sub:
        txt(s, sub, cx0, ct0 + ch0, cw0, Inches(0.3),
            size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)
    if i < len(pa_steps)-1:
        arrow(s, cx0 + cw0, ct0 + ch0 * 0.3, Inches(0.3))
    cx0 += cw0 + Inches(0.3)

txt(s, "作用：生成空间维度权重图（H×W×1），对图像不同区域（浓雾/薄雾/清晰）差异化处理",
    Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.5),
    size=Pt(16), color=ACCENT2, italic=True)

# CA vs PA comparison
txt(s, "CA 与 PA 对比", Inches(0.5), Inches(4.65), Inches(5), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT2)

compare = [
    ("维度", "通道  [C, 1, 1]", "空间  [1, H, W]"),
    ("压缩方式", "全局平均池化", "1×1 卷积"),
    ("输出通道", "C 个权重", "1 个掩码"),
    ("关注点", "哪些特征重要", "哪些位置重要"),
]
headers = ["属性", "CALayer (通道)", "PALayer (像素)"]
tw = [Inches(2.0), Inches(4.5), Inches(4.5)]
tx = [Inches(0.5), Inches(2.6), Inches(7.2)]
ty2 = Inches(5.1)
for col_i, (h, w) in enumerate(zip(headers, tw)):
    box(s, tx[col_i], ty2, w, Inches(0.38), fill=ACCENT)
    txt(s, h, tx[col_i], ty2 + Inches(0.04), w, Inches(0.3),
        size=Pt(13), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
for r, row in enumerate(compare):
    for col_i, (val, w) in enumerate(zip(row, tw)):
        bg_c = BLOCK_BG if r % 2 == 0 else RGBColor(0x22,0x38,0x52)
        box(s, tx[col_i], ty2 + Inches(0.38) + r * Inches(0.38), w, Inches(0.38),
            fill=bg_c, border=RGBColor(0x33,0x55,0x77), border_w=Pt(0.5))
        txt(s, val, tx[col_i] + Inches(0.05), ty2 + Inches(0.42) + r * Inches(0.38),
            w - Inches(0.1), Inches(0.3),
            size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — 特征融合
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "04  特征融合")
txt(s, "跨组特征融合策略", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9),
    fill=BLOCK_BG, border=YELLOW, border_w=Pt(1.5))
txt(s, "out = w₁·G1 + w₂·G2 + w₃·G3    where  [w₁,w₂,w₃] = Sigmoid(FC(Concat([G1,G2,G3])))",
    Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.8),
    size=Pt(17), color=YELLOW, align=PP_ALIGN.CENTER)

# Fusion diagram
g_colors = [GREEN, ACCENT, ORANGE]
g_labels = ["G1 Output\n[B,64,H,W]", "G2 Output\n[B,64,H,W]", "G3 Output\n[B,64,H,W]"]
gx = [Inches(0.6), Inches(2.7), Inches(4.8)]
gy = Inches(2.75)
gw = Inches(1.8)
gh = Inches(0.8)
for lbl, gc, lx in zip(g_labels, g_colors, gx):
    box(s, lx, gy, gw, gh, fill=BLOCK_BG, border=gc, border_w=Pt(1.5))
    txt(s, lbl, lx, gy + Inches(0.05), gw, gh - Inches(0.05),
        size=Pt(12), color=gc, align=PP_ALIGN.CENTER)

# concat box
box(s, Inches(0.6), Inches(3.75), Inches(6.0), Inches(0.75),
    fill=BLOCK_BG, border=ACCENT2, border_w=Pt(1.5))
txt(s, "Concat → [B, 192, H, W]", Inches(0.6), Inches(3.78),
    Inches(6.0), Inches(0.7), size=Pt(14), color=ACCENT2, align=PP_ALIGN.CENTER)

# channel attn box
box(s, Inches(0.6), Inches(4.65), Inches(6.0), Inches(0.75),
    fill=BLOCK_BG, border=YELLOW, border_w=Pt(1.5))
txt(s, "Cross-Group CA → weights [B, 3, 64, 1, 1]",
    Inches(0.6), Inches(4.68), Inches(6.0), Inches(0.7),
    size=Pt(14), color=YELLOW, align=PP_ALIGN.CENTER)

# weighted sum
box(s, Inches(0.6), Inches(5.55), Inches(6.0), Inches(0.75),
    fill=BLOCK_BG, border=GREEN, border_w=Pt(1.5))
txt(s, "Weighted Sum → [B, 64, H, W]",
    Inches(0.6), Inches(5.58), Inches(6.0), Inches(0.7),
    size=Pt(14), color=GREEN, align=PP_ALIGN.CENTER)

# code snippet right side
txt(s, "Fusion 核心代码", Inches(7.2), Inches(2.6), Inches(5), Inches(0.4),
    size=Pt(15), bold=True, color=YELLOW)
fusion_code = (
    "res1 = self.g1(x)\n"
    "res2 = self.g2(res1)\n"
    "res3 = self.g3(res2)\n"
    "# Cross-group channel attention\n"
    "w = self.ca(\n"
    "    torch.cat([res1,res2,res3], dim=1))\n"
    "w = w.view(-1, 3, 64)[:,:,:,None,None]\n"
    "out = (w[:,0]*res1\n"
    "     + w[:,1]*res2\n"
    "     + w[:,2]*res3)\n"
    "out = self.palayer(out)\n"
    "x   = self.post(out)\n"
    "return x + x1  # global residual"
)
box(s, Inches(7.2), Inches(3.05), Inches(5.7), Inches(3.3), fill=RGBColor(0x10,0x20,0x30))
txt(s, fusion_code, Inches(7.35), Inches(3.1), Inches(5.5), Inches(3.2),
    size=Pt(11), color=ACCENT2, italic=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — 训练流程
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "05  训练流程")
txt(s, "训练流程与损失函数", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Training pipeline
pipeline = [
    ("加载\n有雾/清晰图对", ACCENT),
    ("随机裁剪\n& 翻转增广", ACCENT),
    ("归一化\n输入", ACCENT),
    ("FFA-Net\n前向传播", GREEN),
    ("L1 Loss\n(+感知损失)", YELLOW),
    ("反向传播\n更新参数", ORANGE),
    ("余弦学习率\n衰减调度", RED),
]
px = Inches(0.3)
pt = Inches(1.65)
pw = Inches(1.7)
ph = Inches(0.9)
for i, (lbl, col) in enumerate(pipeline):
    box(s, px, pt, pw, ph, fill=BLOCK_BG, border=col, border_w=Pt(1.5))
    txt(s, lbl, px, pt + Inches(0.05), pw, ph - Inches(0.05),
        size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(pipeline)-1:
        arrow(s, px + pw, pt + ph * 0.3, Inches(0.22))
    px += pw + Inches(0.22)

# Loss function details
txt(s, "损失函数详解", Inches(0.5), Inches(2.75), Inches(5), Inches(0.4),
    size=Pt(17), bold=True, color=YELLOW)

box(s, Inches(0.5), Inches(3.2), Inches(5.8), Inches(1.0), fill=BLOCK_BG, border=YELLOW, border_w=Pt(1.5))
txt(s, "L_total = L1(pred, GT)  +  0.04 × L_perceptual(pred, GT)",
    Inches(0.7), Inches(3.3), Inches(5.4), Inches(0.8),
    size=Pt(14), color=YELLOW, align=PP_ALIGN.CENTER)

txt(s, "• L1 Loss：逐像素绝对误差，对噪声鲁棒",
    Inches(0.5), Inches(4.3), Inches(5.8), Inches(0.4), size=Pt(14), color=WHITE)
txt(s, "• Perceptual Loss（可选）：VGG16前16层特征差异，保持语义/纹理",
    Inches(0.5), Inches(4.75), Inches(5.8), Inches(0.4), size=Pt(14), color=WHITE)
txt(s, "• 权重 0.04 防止感知损失主导训练",
    Inches(0.5), Inches(5.2), Inches(5.8), Inches(0.4), size=Pt(14), color=GRAY)

# Learning rate schedule
txt(s, "学习率调度（余弦衰减）", Inches(6.8), Inches(2.75), Inches(5), Inches(0.4),
    size=Pt(17), bold=True, color=ORANGE)
box(s, Inches(6.8), Inches(3.2), Inches(6.0), Inches(0.8), fill=BLOCK_BG, border=ORANGE, border_w=Pt(1.5))
txt(s, "lr(t) = 0.5 × (1 + cos(t·π / T)) × lr₀",
    Inches(7.0), Inches(3.3), Inches(5.6), Inches(0.6),
    size=Pt(16), color=ORANGE, align=PP_ALIGN.CENTER)

hyper = [
    ("优化器", "Adam  (β₁=0.9, β₂=0.999)"),
    ("初始学习率", "1e-4"),
    ("批大小", "2  (ITS) / 2  (OTS)"),
    ("训练步数", "500K (ITS) / 1M (OTS)"),
    ("评估间隔", "每 5000 步"),
    ("裁剪尺寸", "240×240"),
]
txt(s, "超参数配置", Inches(6.8), Inches(4.15), Inches(5), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT2)
hy = Inches(4.58)
for k, v in hyper:
    box(s, Inches(6.8), hy, Inches(2.6), Inches(0.38), fill=BLOCK_BG, border=ACCENT, border_w=Pt(0.5))
    txt(s, k, Inches(6.9), hy + Inches(0.04), Inches(2.4), Inches(0.3),
        size=Pt(13), color=ACCENT)
    box(s, Inches(9.45), hy, Inches(3.35), Inches(0.38), fill=BLOCK_BG, border=ACCENT2, border_w=Pt(0.5))
    txt(s, v, Inches(9.55), hy + Inches(0.04), Inches(3.15), Inches(0.3),
        size=Pt(13), color=WHITE)
    hy += Inches(0.4)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — 数据集
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "06  数据集")
txt(s, "RESIDE 数据集", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

datasets = [
    ("ITS\nIndoor Training Set", "13,990 张有雾图像\n1,399 张清晰图像\n室内合成雾霾\n分辨率多样", ACCENT, "训练集"),
    ("OTS\nOutdoor Training Set", "313,950 张有雾图像\n8,970 张清晰图像\n室外合成雾霾\n大规模训练", GREEN, "训练集"),
    ("SOTS Indoor\nTest Set", "500 对室内测试图\n合成雾霾\n评估 PSNR / SSIM", ORANGE, "测试集"),
    ("SOTS Outdoor\nTest Set", "500 对室外测试图\n合成雾霾\n评估 PSNR / SSIM", YELLOW, "测试集"),
]
dx = Inches(0.4)
dt = Inches(1.6)
dw = Inches(3.05)
dh = Inches(3.2)
for title, info, col, tag in datasets:
    box(s, dx, dt, dw, dh, fill=BLOCK_BG, border=col, border_w=Pt(1.5))
    box(s, dx, dt, dw, Inches(0.35), fill=col)
    txt(s, tag, dx + Inches(0.05), dt + Inches(0.04), Inches(0.8), Inches(0.28),
        size=Pt(11), bold=True, color=BG_DARK)
    txt(s, title, dx + Inches(0.1), dt + Inches(0.42), dw - Inches(0.2), Inches(0.85),
        size=Pt(14), bold=True, color=col)
    txt(s, info, dx + Inches(0.1), dt + Inches(1.3), dw - Inches(0.2), Inches(1.7),
        size=Pt(13), color=WHITE)
    dx += dw + Inches(0.12)

# Data augmentation
txt(s, "数据增强策略", Inches(0.5), Inches(4.95), Inches(5), Inches(0.4),
    size=Pt(17), bold=True, color=ACCENT)
aug_items = [
    "随机水平翻转（概率 0.5）",
    "随机旋转 0° / 90° / 180° / 270°",
    "随机裁剪 240×240 patch",
    "归一化：mean=[0.64, 0.60, 0.58], std=[0.14, 0.15, 0.152]",
]
bullet_list(s, aug_items, Inches(0.5), Inches(5.4), Inches(6.0), Inches(2.0),
            size=Pt(14), dot_color=ACCENT)

# naming convention
txt(s, "文件命名规则", Inches(7.0), Inches(4.95), Inches(5), Inches(0.4),
    size=Pt(17), bold=True, color=ACCENT2)
txt(s, "有雾图像：  <id>_<beta>_<A>.png\n清晰图像：  <id>.png\n\n例：1400_0.9_0.16.jpg  →  清晰图 1400.png",
    Inches(7.0), Inches(5.4), Inches(6.0), Inches(1.8),
    size=Pt(14), color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — 实验结果
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "07  实验结果")
txt(s, "性能评估与对比", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Table
headers = ["方法", "室内 PSNR (dB)", "室内 SSIM", "室外 PSNR (dB)", "室外 SSIM"]
rows = [
    ("DCP",       "16.62", "0.8179", "19.13", "0.8148"),
    ("AOD-Net",   "19.06", "0.8504", "20.29", "0.8765"),
    ("DehazeNet", "21.14", "0.8472", "22.46", "0.8514"),
    ("GFN",       "22.30", "0.8800", "21.55", "0.8444"),
    ("GCANet",    "30.23", "0.9800", "—",     "—"),
    ("FFA-Net\n(Ours)", "36.39", "0.9886", "33.57", "0.9840"),
]
col_w2 = [Inches(2.4), Inches(2.2), Inches(2.0), Inches(2.2), Inches(2.0)]
col_x = [Inches(0.5)]
for w in col_w2[:-1]:
    col_x.append(col_x[-1] + w + Inches(0.05))
th = Inches(0.45)
ty3 = Inches(1.65)

for ci, (h, w) in enumerate(zip(headers, col_w2)):
    box(s, col_x[ci], ty3, w, th, fill=ACCENT)
    txt(s, h, col_x[ci], ty3 + Inches(0.04), w, th - Inches(0.04),
        size=Pt(13), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    is_ours = row[0].startswith("FFA")
    bg_c = RGBColor(0x06,0x3A,0x1F) if is_ours else (BLOCK_BG if ri % 2 == 0 else RGBColor(0x1F,0x35,0x50))
    border_c = GREEN if is_ours else RGBColor(0x33,0x55,0x77)
    rh = Inches(0.55) if is_ours else Inches(0.43)
    for ci, (val, w) in enumerate(zip(row, col_w2)):
        box(s, col_x[ci], ty3 + th + ri * Inches(0.43), w, rh,
            fill=bg_c, border=border_c, border_w=Pt(1.5 if is_ours else 0.5))
        c_ = GREEN if is_ours else WHITE
        txt(s, val, col_x[ci] + Inches(0.05), ty3 + th + ri * Inches(0.43) + Inches(0.04),
            w - Inches(0.1), rh - Inches(0.04),
            size=Pt(13), bold=is_ours, color=c_, align=PP_ALIGN.CENTER)

# gain analysis
txt(s, "相对 GCANet（次优方法）的提升", Inches(0.5), Inches(5.1), Inches(6), Inches(0.4),
    size=Pt(15), bold=True, color=GREEN)
gains = [
    ("室内 PSNR", "+6.16 dB", GREEN),
    ("室内 SSIM", "+0.0086", GREEN),
    ("室外 PSNR", "+14.44 dB（GCANet 无室外结果）", YELLOW),
]
gy2 = Inches(5.55)
for k, v, col in gains:
    box(s, Inches(0.5), gy2, Inches(2.2), Inches(0.38), fill=BLOCK_BG, border=ACCENT, border_w=Pt(0.5))
    txt(s, k, Inches(0.6), gy2 + Inches(0.04), Inches(2.0), Inches(0.3),
        size=Pt(13), color=ACCENT)
    box(s, Inches(2.75), gy2, Inches(5.0), Inches(0.38), fill=BLOCK_BG, border=col, border_w=Pt(0.5))
    txt(s, v, Inches(2.85), gy2 + Inches(0.04), Inches(4.8), Inches(0.3),
        size=Pt(13), bold=True, color=col)
    gy2 += Inches(0.42)

# metric explanation
txt(s, "指标说明", Inches(8.2), Inches(5.1), Inches(4.5), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT2)
txt(s, "PSNR（峰值信噪比）：越高越好，单位 dB\nPSNR = 20·log₁₀(1 / RMSE)\n\nSSIM（结构相似度）：范围 [0,1]，越接近1越好\n综合亮度/对比度/结构三方面评估",
    Inches(8.2), Inches(5.55), Inches(4.8), Inches(1.7),
    size=Pt(13), color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — 效果展示
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "08  效果展示")
txt(s, "去雾效果可视化", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Row 1 - indoor sample
txt(s, "室内样本（ITS 数据集）", Inches(0.5), Inches(1.55), Inches(5), Inches(0.38),
    size=Pt(14), bold=True, color=ACCENT)
add_image(s, "/home/user/FFA-Net/fig/1400_2.png",
          Inches(0.4), Inches(1.95), Inches(6.1), Inches(2.4))
add_image(s, "/home/user/FFA-Net/fig/1400_2_FFA.png",
          Inches(6.8), Inches(1.95), Inches(6.1), Inches(2.4))

txt(s, "有雾输入", Inches(2.5), Inches(4.4), Inches(2), Inches(0.4),
    size=Pt(14), color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "FFA-Net 输出", Inches(9.1), Inches(4.4), Inches(2), Inches(0.4),
    size=Pt(14), color=GREEN, align=PP_ALIGN.CENTER)

# Row 2 - outdoor sample
txt(s, "室外样本（OTS 数据集）", Inches(0.5), Inches(4.85), Inches(5), Inches(0.38),
    size=Pt(14), bold=True, color=ACCENT)
add_image(s, "/home/user/FFA-Net/fig/0099_0.9_0.16.jpg",
          Inches(0.4), Inches(5.25), Inches(6.1), Inches(1.9))
add_image(s, "/home/user/FFA-Net/fig/0099_0_FFA.png",
          Inches(6.8), Inches(5.25), Inches(6.1), Inches(1.9))

txt(s, "有雾输入", Inches(2.5), Inches(7.2), Inches(2), Inches(0.35),
    size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "FFA-Net 输出", Inches(9.1), Inches(7.2), Inches(2), Inches(0.35),
    size=Pt(13), color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — 代码结构
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "09  代码结构")
txt(s, "项目代码结构", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

# Directory tree
tree = (
    "FFA-Net/\n"
    "├── net/\n"
    "│   ├── main.py          # 训练入口：训练循环、eval、checkpoint\n"
    "│   ├── test.py          # 推理入口：加载模型对图像去雾\n"
    "│   ├── option.py        # 超参数 CLI 定义\n"
    "│   ├── data_utils.py    # RESIDE 数据集加载、增强\n"
    "│   ├── metrics.py       # PSNR / SSIM 实现\n"
    "│   └── models/\n"
    "│       ├── FFA.py       # 核心网络：PALayer/CALayer/Block/Group/FFA\n"
    "│       └── PerceptualLoss.py  # VGG16 感知损失\n"
    "├── fig/                 # 示例图片\n"
    "└── README.md"
)
box(s, Inches(0.5), Inches(1.6), Inches(7.5), Inches(4.8), fill=RGBColor(0x08,0x14,0x22))
txt(s, tree, Inches(0.65), Inches(1.7), Inches(7.2), Inches(4.6),
    size=Pt(13), color=ACCENT2, italic=True)

# Key files detail
files_detail = [
    ("main.py", ["lr_schedule_cosdecay() → 余弦衰减",
                 "train() → 训练主循环，支持断点续训",
                 "test() → 评估 PSNR/SSIM",
                 "支持 TensorboardX 可视化（注释状态）"]),
    ("data_utils.py", ["RESIDE_Dataset → 自定义 Dataset",
                        "augData() → 翻转 + 旋转 + 归一化",
                        "4 个预建 DataLoader（ITS/OTS × train/test）"]),
    ("option.py", ["argparse CLI 参数定义",
                   "自动检测 CUDA，自动创建目录",
                   "model_name 由参数组合生成"]),
]
fd_x = Inches(8.2)
fd_y = Inches(1.6)
fd_w = Inches(4.9)
fd_h = Inches(1.5)
for fname, details in files_detail:
    box(s, fd_x, fd_y, fd_w, fd_h, fill=BLOCK_BG, border=ACCENT, border_w=Pt(1))
    txt(s, fname, fd_x + Inches(0.1), fd_y + Inches(0.05), fd_w - Inches(0.2), Inches(0.38),
        size=Pt(14), bold=True, color=ACCENT)
    dy = fd_y + Inches(0.45)
    for d in details:
        txt(s, "· " + d, fd_x + Inches(0.1), dy, fd_w - Inches(0.2), Inches(0.3),
            size=Pt(11), color=GRAY)
        dy += Inches(0.3)
    fd_y += fd_h + Inches(0.15)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — 总结与展望
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
section_badge(s, "10  总结展望")
txt(s, "总结与未来展望", Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(1.45))

txt(s, "核心贡献", Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
    size=Pt(18), bold=True, color=GREEN)
contributions = [
    "提出像素级注意力（PALayer），差异化处理不同浓度雾区域",
    "提出通道级注意力（CALayer），对不同颜色/语义通道加权",
    "设计跨组特征融合（FFA），动态整合多尺度特征",
    "室内 PSNR 36.39 dB / SSIM 0.9886，大幅领先现有 SOTA",
    "代码简洁，PyTorch 实现，易于复现与扩展",
]
bullet_list(s, contributions, Inches(0.5), Inches(2.1), Inches(5.8), Inches(2.8),
            size=Pt(14), dot_color=GREEN)

box(s, Inches(6.3), Inches(1.55), Pt(2), Inches(5.5), fill=BLOCK_BG)

txt(s, "局限性", Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.45),
    size=Pt(18), bold=True, color=ORANGE)
limitations = [
    "仅在合成数据集上验证，真实世界泛化性有待提升",
    "分辨率受限于 GPU 显存，高分辨率图像需分块处理",
    "固定 gps=3，Group 数量不可灵活配置",
    "推理速度未针对边缘设备优化",
]
bullet_list(s, limitations, Inches(6.8), Inches(2.1), Inches(5.8), Inches(2.2),
            size=Pt(14), dot_color=ORANGE)

txt(s, "未来方向", Inches(6.8), Inches(4.4), Inches(5.8), Inches(0.45),
    size=Pt(18), bold=True, color=ACCENT)
futures = [
    "引入 Transformer / 自注意力机制增强全局感知",
    "结合对比学习进行无监督去雾",
    "扩展至视频去雾，增加时序一致性约束",
    "轻量化模型以支持移动端部署",
]
bullet_list(s, futures, Inches(6.8), Inches(4.9), Inches(5.8), Inches(2.2),
            size=Pt(14), dot_color=ACCENT)

# bottom performance summary strip
box(s, Inches(0.5), Inches(5.55), Inches(5.3), Inches(1.5), fill=BLOCK_BG, border=GREEN, border_w=Pt(1.5))
txt(s, "最终性能", Inches(0.7), Inches(5.65), Inches(4.8), Inches(0.38),
    size=Pt(14), bold=True, color=GREEN)
txt(s, "室内：PSNR 36.39 dB  /  SSIM 0.9886\n室外：PSNR 33.57 dB  /  SSIM 0.9840",
    Inches(0.7), Inches(6.05), Inches(4.8), Inches(0.9),
    size=Pt(15), bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, prs.slide_width, Inches(0.08), fill=ACCENT)
box(s, 0, prs.slide_height - Inches(0.08), prs.slide_width, Inches(0.08), fill=ACCENT)

c = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x00,0x5F,0x7A)
c.line.fill.background()

txt(s, "谢谢！", Inches(1), Inches(1.8), Inches(11), Inches(1.5),
    size=Pt(80), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "FFA-Net: Feature Fusion Attention Network for Single Image Dehazing",
    Inches(1), Inches(3.5), Inches(11), Inches(0.7),
    size=Pt(20), color=ACCENT2, align=PP_ALIGN.CENTER)
txt(s, "AAAI 2020  ·  arXiv:1911.07559",
    Inches(1), Inches(4.3), Inches(11), Inches(0.5),
    size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

txt(s, "Q & A",
    Inches(1), Inches(5.2), Inches(11), Inches(0.8),
    size=Pt(36), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════
output_path = "/home/user/FFA-Net/FFA-Net_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
